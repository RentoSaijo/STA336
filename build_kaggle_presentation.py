from __future__ import annotations

import csv
import shutil
from collections import Counter, defaultdict
from pathlib import Path
from typing import Iterable

import matplotlib.pyplot as plt
from matplotlib import font_manager
from matplotlib.patches import FancyBboxPatch
from matplotlib.ticker import FuncFormatter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    from fontTools.ttLib import TTCollection
except ImportError:  # pragma: no cover - defensive fallback
    TTCollection = None


ROOT = Path(__file__).resolve().parent
PPTX_PATH = ROOT / "kaggle.pptx"
BACKUP_PATH = ROOT / "kaggle_theme_backup.pptx"
ASSET_DIR = ROOT / "generated" / "kaggle_ppt_assets"
FONT = "Cambria Math"
OFFICE_CAMBRIA_TTC = Path("/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts/Cambria.ttc")

WHITE = RGBColor(245, 245, 245)
MUTED = RGBColor(188, 188, 188)
PANEL = RGBColor(60, 60, 60)
PANEL_DARK = RGBColor(52, 52, 52)
ACCENT_BLUE = RGBColor(15, 158, 213)
ACCENT_GREEN = RGBColor(78, 167, 46)
ACCENT_ORANGE = RGBColor(233, 113, 50)
ACCENT_TEAL = RGBColor(21, 96, 130)


def load_rows(path: Path) -> list[dict[str, str]]:
    with path.open(newline="") as f:
        return list(csv.DictReader(f))


def compute_stats(train_rows: list[dict[str, str]], test_rows: list[dict[str, str]]) -> dict:
    stats: dict[str, object] = {}

    train_n = len(train_rows)
    test_n = len(test_rows)
    true_n = sum(r["Transported"] == "True" for r in train_rows)
    false_n = sum(r["Transported"] == "False" for r in train_rows)
    majority_acc = max(true_n, false_n) / train_n

    stats["train_n"] = train_n
    stats["test_n"] = test_n
    stats["true_n"] = true_n
    stats["false_n"] = false_n
    stats["true_pct"] = true_n / train_n
    stats["false_pct"] = false_n / train_n
    stats["majority_acc"] = majority_acc

    missing_train = {}
    for col in train_rows[0]:
        missing_train[col] = sum(r[col] == "" for r in train_rows)
    stats["missing_train"] = missing_train

    missing_test = {}
    for col in test_rows[0]:
        missing_test[col] = sum(r[col] == "" for r in test_rows)
    stats["missing_test"] = missing_test

    row_missing_counts = []
    cols = list(train_rows[0].keys())
    for r in train_rows:
        row_missing_counts.append(sum(r[c] == "" for c in cols))
    stats["complete_rows"] = row_missing_counts.count(0)
    stats["one_missing_rows"] = row_missing_counts.count(1)
    stats["two_missing_rows"] = row_missing_counts.count(2)
    stats["three_missing_rows"] = row_missing_counts.count(3)

    for r in train_rows:
        spends = [float(r[c]) if r[c] else 0.0 for c in ["RoomService", "FoodCourt", "ShoppingMall", "Spa", "VRDeck"]]
        r["TotalSpend"] = sum(spends)
        r["NoSpend"] = r["TotalSpend"] == 0
        if r["Cabin"]:
            deck, _, side = r["Cabin"].split("/")
        else:
            deck, side = "Missing", "Missing"
        r["Deck"] = deck
        r["Side"] = side
        gid, pos = r["PassengerId"].split("_")
        r["GroupID"] = gid
        r["GroupPos"] = int(pos)

    by_group: dict[str, list[dict[str, str]]] = defaultdict(list)
    for r in train_rows:
        by_group[r["GroupID"]].append(r)
    for group in by_group.values():
        group_size = max(r["GroupPos"] for r in group)
        for r in group:
            r["GroupSize"] = group_size

    def rate(subset: Iterable[dict[str, str]]) -> float:
        subset = list(subset)
        return sum(r["Transported"] == "True" for r in subset) / len(subset)

    stats["rate_cryo_true"] = rate(r for r in train_rows if r["CryoSleep"] == "True")
    stats["rate_cryo_false"] = rate(r for r in train_rows if r["CryoSleep"] == "False")
    stats["rate_nospend"] = rate(r for r in train_rows if r["NoSpend"])
    stats["rate_posspend"] = rate(r for r in train_rows if not r["NoSpend"])
    stats["rate_side_p"] = rate(r for r in train_rows if r["Side"] == "P")
    stats["rate_side_s"] = rate(r for r in train_rows if r["Side"] == "S")

    age_le_12 = [r for r in train_rows if r["Age"] and float(r["Age"]) <= 12]
    age_gt_12 = [r for r in train_rows if r["Age"] and float(r["Age"]) > 12]
    stats["rate_child"] = rate(age_le_12)
    stats["rate_adult"] = rate(age_gt_12)

    size_rates = {}
    for size in sorted({r["GroupSize"] for r in train_rows}):
        subset = [r for r in train_rows if r["GroupSize"] == size]
        size_rates[size] = rate(subset)
    stats["group_size_rates"] = size_rates

    deck_rates = {}
    for deck in sorted({r["Deck"] for r in train_rows if r["Deck"] != "Missing"}):
        subset = [r for r in train_rows if r["Deck"] == deck]
        deck_rates[deck] = rate(subset)
    stats["deck_rates"] = deck_rates

    deck_home = defaultdict(Counter)
    for r in train_rows:
        if r["Deck"] != "Missing" and r["HomePlanet"]:
            deck_home[r["Deck"]][r["HomePlanet"]] += 1
    stats["deck_home"] = {deck: dict(counter) for deck, counter in deck_home.items()}

    def share_same(key_fn) -> tuple[float, int]:
        flags = []
        for group in by_group.values():
            obs = [key_fn(r) for r in group if key_fn(r) not in ("", "Missing", None)]
            if len(obs) <= 1:
                continue
            flags.append(len(set(obs)) == 1)
        return sum(flags) / len(flags), len(flags)

    stats["same_homeplanet"], stats["same_homeplanet_n"] = share_same(lambda r: r["HomePlanet"])
    stats["same_destination"], stats["same_destination_n"] = share_same(lambda r: r["Destination"])
    stats["same_cabinside"], stats["same_cabinside_n"] = share_same(lambda r: r["Side"])
    stats["same_cabindeck"], stats["same_cabindeck_n"] = share_same(lambda r: r["Deck"])
    stats["uniform_group_outcome"] = sum(
        len({r["Transported"] for r in group}) == 1 for group in by_group.values()
    ) / len(by_group)

    cryo_true = [r for r in train_rows if r["CryoSleep"] == "True"]
    cryo_false = [r for r in train_rows if r["CryoSleep"] == "False"]
    stats["zero_spend_share_cryo_true"] = sum(r["NoSpend"] for r in cryo_true) / len(cryo_true)
    stats["zero_spend_share_cryo_false"] = sum(r["NoSpend"] for r in cryo_false) / len(cryo_false)

    zero = [r for r in train_rows if r["NoSpend"]]
    positive = [r for r in train_rows if not r["NoSpend"]]
    stats["zero_cryo_true_share"] = sum(r["CryoSleep"] == "True" for r in zero) / len(zero)
    stats["zero_cryo_false_share"] = sum(r["CryoSleep"] == "False" for r in zero) / len(zero)
    stats["positive_cryo_true_share"] = sum(r["CryoSleep"] == "True" for r in positive) / len(positive)

    stats["starter_logistic"] = 0.6266427718
    stats["engineered_logistic"] = 0.7951
    stats["tuned_lightgbm"] = 0.8066274343
    stats["patrick_rf_oob"] = 0.8060508455
    stats["patrick_hgb_cv"] = 0.8077817018
    stats["majority_error_reduction"] = 1 - ((1 - stats["tuned_lightgbm"]) / (1 - stats["majority_acc"]))

    stats["maria_cv_models"] = [
        ("Logistic", 0.7951),
        ("CatBoost", 0.7955),
        ("AdaBoost", 0.7952),
        ("LightGBM", 0.7985),
        ("Random Forest", 0.7997),
        ("SVC", 0.8005),
        ("XGBoost", 0.8020),
        ("Patrick RF OOB", 0.8061),
        ("Patrick HistGB CV", 0.8078),
        ("Tuned LightGBM", 0.8066),
    ]

    return stats


def pct(x: float) -> str:
    return f"{100 * x:.2f}%"


def fmt_rate(x: float) -> str:
    return f"{100 * x:.1f}%"


def ensure_chart_font() -> str:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    available = {f.name for f in font_manager.fontManager.ttflist}
    if FONT in available:
        return FONT

    if OFFICE_CAMBRIA_TTC.exists() and TTCollection is not None:
        extracted = ASSET_DIR / "cambria_math_extracted.ttf"
        if not extracted.exists():
            ttc = TTCollection(str(OFFICE_CAMBRIA_TTC))
            ttc.fonts[1].save(str(extracted))
        font_manager.fontManager.addfont(str(extracted))
        available = {f.name for f in font_manager.fontManager.ttflist}
        if FONT in available:
            return FONT

    if OFFICE_CAMBRIA_TTC.exists():
        font_manager.fontManager.addfont(str(OFFICE_CAMBRIA_TTC))
        available = {f.name for f in font_manager.fontManager.ttflist}
        if "Cambria" in available:
            return "Cambria"

    return "DejaVu Serif"


def make_charts(stats: dict) -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    chart_font = ensure_chart_font()
    plt.rcParams.update(
        {
            "font.family": chart_font,
            "font.size": 10,
            "figure.facecolor": "none",
            "axes.facecolor": "none",
            "savefig.facecolor": "none",
            "text.color": "#F5F5F5",
            "axes.labelcolor": "#F5F5F5",
            "xtick.color": "#E0E0E0",
            "ytick.color": "#E0E0E0",
            "axes.edgecolor": "#BDBDBD",
        }
    )

    def percent_formatter(val, _):
        return f"{100 * val:.0f}%"

    chart_paths: dict[str, Path] = {}

    # Missingness chart
    missing_items = [
        (k, v / stats["train_n"])
        for k, v in stats["missing_train"].items()
        if v > 0 and k not in {"Transported", "PassengerId"}
    ]
    missing_items.sort(key=lambda kv: kv[1])
    fig, ax = plt.subplots(figsize=(7.2, 4.2))
    labels = [k for k, _ in missing_items]
    vals = [v for _, v in missing_items]
    colors = ["#0F9ED5" if k in {"CryoSleep", "Cabin", "HomePlanet"} else "#4EA72E" for k in labels]
    ax.barh(labels, vals, color=colors)
    ax.set_title("Training-Set Missingness by Feature", pad=12, fontsize=13)
    ax.xaxis.set_major_formatter(FuncFormatter(percent_formatter))
    ax.set_xlabel("Share of rows with blanks")
    ax.grid(axis="x", alpha=0.18)
    for i, val in enumerate(vals):
        ax.text(val + 0.0008, i, f"{100*val:.2f}%", va="center", fontsize=9)
    fig.tight_layout()
    chart_paths["missingness"] = ASSET_DIR / "missingness.png"
    fig.savefig(chart_paths["missingness"], dpi=180, transparent=True)
    plt.close(fig)

    # Validation ladder chart
    ladder = [
        ("Majority baseline", stats["majority_acc"]),
        ("Spa-only logistic", stats["starter_logistic"]),
        ("Engineered logistic", stats["engineered_logistic"]),
        ("Tuned LightGBM", stats["tuned_lightgbm"]),
        ("Patrick HistGB CV", stats["patrick_hgb_cv"]),
    ]
    fig, ax = plt.subplots(figsize=(6.8, 4.0))
    labels = [k for k, _ in ladder]
    vals = [v for _, v in ladder]
    colors = ["#8A8A8A", "#E97132", "#156082", "#4EA72E", "#0F9ED5"]
    ax.bar(labels, vals, color=colors)
    ax.set_ylim(0.45, 0.84)
    ax.yaxis.set_major_formatter(FuncFormatter(percent_formatter))
    ax.set_title("Accuracy Gains from Baseline to Strong Public Models", pad=12, fontsize=13)
    ax.grid(axis="y", alpha=0.18)
    plt.setp(ax.get_xticklabels(), rotation=18, ha="right")
    for idx, val in enumerate(vals):
        ax.text(idx, val + 0.008, f"{100*val:.1f}%", ha="center", fontsize=9)
    fig.tight_layout()
    chart_paths["validation"] = ASSET_DIR / "validation_ladder.png"
    fig.savefig(chart_paths["validation"], dpi=180, transparent=True)
    plt.close(fig)

    # Feature signals chart
    fig, axs = plt.subplots(2, 2, figsize=(8.2, 5.8))
    cryo_labels = ["Cryo=False", "Cryo=True"]
    cryo_vals = [stats["rate_cryo_false"], stats["rate_cryo_true"]]
    axs[0, 0].bar(cryo_labels, cryo_vals, color=["#E97132", "#4EA72E"])
    axs[0, 0].set_title("Transport Rate by CryoSleep")
    axs[0, 0].set_ylim(0, 1)

    spend_labels = ["Positive spend", "Zero spend"]
    spend_vals = [stats["rate_posspend"], stats["rate_nospend"]]
    axs[0, 1].bar(spend_labels, spend_vals, color=["#E97132", "#4EA72E"])
    axs[0, 1].set_title("Transport Rate by Total Spend")
    axs[0, 1].set_ylim(0, 1)

    side_labels = ["Port", "Starboard"]
    side_vals = [stats["rate_side_p"], stats["rate_side_s"]]
    axs[1, 0].bar(side_labels, side_vals, color=["#156082", "#0F9ED5"])
    axs[1, 0].set_title("Transport Rate by Ship Side")
    axs[1, 0].set_ylim(0, 1)

    deck_labels = list(stats["deck_rates"].keys())
    deck_vals = [stats["deck_rates"][d] for d in deck_labels]
    axs[1, 1].bar(deck_labels, deck_vals, color=["#4EA72E", "#0F9ED5", "#156082", "#E97132", "#8A8A8A", "#4EA72E", "#0F9ED5", "#E97132"][: len(deck_labels)])
    axs[1, 1].set_title("Transport Rate by Deck")
    axs[1, 1].set_ylim(0, 1)

    for ax in axs.flat:
        ax.yaxis.set_major_formatter(FuncFormatter(percent_formatter))
        ax.grid(axis="y", alpha=0.18)
        for tick in ax.get_xticklabels():
            tick.set_rotation(12)
    fig.tight_layout()
    chart_paths["signals"] = ASSET_DIR / "feature_signals.png"
    fig.savefig(chart_paths["signals"], dpi=180, transparent=True)
    plt.close(fig)

    # Model comparison chart
    model_items = sorted(stats["maria_cv_models"], key=lambda kv: kv[1])
    fig, ax = plt.subplots(figsize=(7.4, 4.8))
    labels = [k for k, _ in model_items]
    vals = [v for _, v in model_items]
    colors = []
    for label in labels:
        if "Logistic" in label:
            colors.append("#156082")
        elif "LightGBM" in label or "HistGB" in label:
            colors.append("#4EA72E")
        elif "Patrick RF" in label or "Random Forest" in label:
            colors.append("#0F9ED5")
        else:
            colors.append("#E97132")
    ax.barh(labels, vals, color=colors)
    ax.set_xlim(0.785, 0.812)
    ax.xaxis.set_major_formatter(FuncFormatter(percent_formatter))
    ax.set_title("Public Validation Results After Strong Preprocessing", pad=12, fontsize=13)
    ax.grid(axis="x", alpha=0.18)
    for idx, val in enumerate(vals):
        ax.text(val + 0.0004, idx, f"{100*val:.2f}%", va="center", fontsize=9)
    fig.tight_layout()
    chart_paths["models"] = ASSET_DIR / "model_compare.png"
    fig.savefig(chart_paths["models"], dpi=180, transparent=True)
    plt.close(fig)

    # HistGradientBoosting explainer diagram
    fig, ax = plt.subplots(figsize=(8.0, 4.8))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7)
    ax.axis("off")

    def diagram_box(x, y, w, h, title, body, edge):
        patch = FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.03,rounding_size=0.12",
            linewidth=1.6,
            edgecolor=edge,
            facecolor="#3C3C3C",
        )
        ax.add_patch(patch)
        ax.text(x + 0.18, y + h - 0.24, title, color=edge, fontsize=12.5, fontweight="bold", va="top")
        ax.text(x + 0.18, y + h - 0.65, body, color="#F5F5F5", fontsize=10.3, va="top", wrap=True)

    diagram_box(
        0.35,
        4.45,
        2.55,
        1.65,
        "1. Bin Features",
        "Each numeric feature is placed into histogram bins instead of searching every unique split value.",
        "#0F9ED5",
    )
    diagram_box(
        3.25,
        4.45,
        2.55,
        1.65,
        "2. Fit Weak Tree",
        "At round m, fit a tree to the negative gradient of binary log-loss.",
        "#4EA72E",
    )
    diagram_box(
        6.15,
        4.45,
        3.2,
        1.65,
        "3. Update Ensemble",
        "Add the new tree with a small step: F_m(x) = F_(m-1)(x) + nu h_m(x).",
        "#E97132",
    )
    diagram_box(
        1.3,
        1.2,
        7.45,
        1.95,
        "Why it fits Spaceship Titanic",
        "Mixed categorical and numeric features, skewed spending, and strong interactions make stage-wise boosted trees a natural fit.\n"
        "Patrick's notebook uses a categorical-feature mask, no one-hot encoding, and log-transformed expenses.",
        "#156082",
    )

    ax.annotate("", xy=(3.13, 5.28), xytext=(2.93, 5.28), arrowprops=dict(arrowstyle="->", lw=2.0, color="#0F9ED5"))
    ax.annotate("", xy=(6.03, 5.28), xytext=(5.83, 5.28), arrowprops=dict(arrowstyle="->", lw=2.0, color="#4EA72E"))
    ax.annotate("", xy=(5.0, 3.2), xytext=(5.0, 4.35), arrowprops=dict(arrowstyle="->", lw=2.0, color="#F5F5F5"))
    ax.text(5.0, 0.35, "PatrickSVM 10-fold CV mean = 0.80778", ha="center", color="#F5F5F5", fontsize=11.5)
    fig.tight_layout()
    chart_paths["histgb"] = ASSET_DIR / "histgb_diagram.png"
    fig.savefig(chart_paths["histgb"], dpi=180, transparent=True)
    plt.close(fig)

    return chart_paths


def set_run_style(run, size: int, color: RGBColor = WHITE, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def style_text_frame(text_frame, size: int, color: RGBColor = WHITE, bold_first: bool = False) -> None:
    for idx, paragraph in enumerate(text_frame.paragraphs):
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            set_run_style(run, size=size, color=color, bold=bold_first and idx == 0)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size: int = 18,
    color: RGBColor = WHITE,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            set_run_style(run, size=size, color=color, bold=bold)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)


def add_bullets_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: list[str],
    accent: RGBColor,
    title_size: int = 18,
    bullet_size: int = 14,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.2)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        set_run_style(run, title_size, accent, bold=True)

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            set_run_style(run, bullet_size, WHITE, bold=False)


def add_diagram_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: str,
    accent: RGBColor,
    title_size: int = 14,
    body_size: int = 11,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.4)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        set_run_style(run, title_size, accent, bold=True)

    p = tf.add_paragraph()
    p.text = body
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        set_run_style(run, body_size, WHITE, bold=False)


def add_arrow(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    direction: str = "right",
    color: RGBColor = WHITE,
) -> None:
    shape_type = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, 0.62, 0.36, 12.1, 0.55, title, size=27, color=WHITE, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.94), Inches(2.3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.64, 0.99, 11.8, 0.38, subtitle, size=12, color=MUTED)


def add_footer(slide, footer: str) -> None:
    add_textbox(slide, 0.6, 7.08, 12.0, 0.22, footer, size=9, color=MUTED)


def set_notes(slide, text: str) -> None:
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text
    for paragraph in notes_tf.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            set_run_style(run, size=12, color=RGBColor(0, 0, 0), bold=False)


def build_presentation() -> None:
    if not BACKUP_PATH.exists():
        shutil.copy2(PPTX_PATH, BACKUP_PATH)

    train_rows = load_rows(Path("/Users/rsai_91/Downloads/spaceship-titanic/train.csv"))
    test_rows = load_rows(Path("/Users/rsai_91/Downloads/spaceship-titanic/test.csv"))
    stats = compute_stats(train_rows, test_rows)
    charts = make_charts(stats)

    prs = Presentation(BACKUP_PATH)

    # Reuse the three existing slides, then add the rest as blanks.
    while len(prs.slides) < 12:
        prs.slides.add_slide(prs.slide_layouts[6])

    slides = list(prs.slides)

    # Slide 1: title
    title_slide = slides[0]
    clear_slide(title_slide)
    add_textbox(
        title_slide,
        2.4,
        2.35,
        7.2,
        0.7,
        "Kaggle Spring",
        size=28,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        title_slide,
        2.0,
        3.45,
        9.4,
        0.55,
        "Spaceship Titanic: context, public workflows, and ISLR takeaways",
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        title_slide,
        3.2,
        4.15,
        6.9,
        0.35,
        "Built from local CSV analysis and public notebook evidence",
        size=11,
        color=RGBColor(210, 210, 210),
        align=PP_ALIGN.CENTER,
    )
    set_notes(
        title_slide,
        "Open by framing the project as a research synthesis rather than a custom modeling exercise.\n"
        "- Our goal is to explain how strong public Spaceship Titanic workflows work and why they work.\n"
        "- The deck is built from two evidence streams: the attached train/test files and detailed public notebooks and repos.\n"
        "- I will move from context, to data structure, to validation, to feature engineering, to model comparison, and then end with ISLR connections and best practices.\n"
        "- Because the competition uses a rolling leaderboard, the presentation focuses on strong public workflows rather than claiming a single final winning solution."
    )

    # Slide 2: goal and context
    slide = slides[1]
    clear_slide(slide)
    add_title(slide, "Competition Goal and Context", "Why this is a good case study for statistical learning")
    add_bullets_box(
        slide,
        0.7,
        1.3,
        5.8,
        3.2,
        "What Kaggle is asking",
        [
            "Binary classification: predict whether each passenger was Transported.",
            "One hidden-label test set; one True/False prediction per PassengerId.",
            "Ongoing rolling leaderboard, so workflow matters more than any one rank snapshot.",
        ],
        ACCENT_BLUE,
    )
    add_bullets_box(
        slide,
        6.8,
        1.3,
        5.8,
        3.2,
        "Why the problem is nontrivial",
        [
            f"Majority-class baseline = {pct(stats['majority_acc'])}.",
            f"Strong public workflows reach about {pct(stats['tuned_lightgbm'])} to {pct(stats['patrick_hgb_cv'])}.",
            f"That is about a {100*stats['majority_error_reduction']:.1f}% reduction in error rate versus the naive rule.",
        ],
        ACCENT_GREEN,
    )
    add_textbox(
        slide,
        0.9,
        4.95,
        11.7,
        1.25,
        "Accuracy here is reasonable because the target is almost perfectly balanced in the training set: "
        f"{stats['true_n']} True ({pct(stats['true_pct'])}) versus {stats['false_n']} False ({pct(stats['false_pct'])}).",
        size=17,
    )
    add_footer(slide, "Sources: class competition brief; official Kaggle page; local train.csv/test.csv")
    set_notes(
        slide,
        "This slide should establish the problem before any methods appear.\n"
        "- Spaceship Titanic is a classification task, not regression: the target is a Bernoulli outcome, Transported or not.\n"
        "- The hidden test set means participants must validate locally before submission.\n"
        "- The class prompt notes that the competition uses a rolling leaderboard, so there is no single fixed winning solution to memorize.\n"
        "- The target balance matters: 4,378 True versus 4,315 False in the attached training file. That is why accuracy is a defensible metric here.\n"
        "- A nice quantitative line to say aloud is that the majority rule gives only 50.36% accuracy, whereas strong public pipelines are around 80.7%, which cuts error by about 61%."
    )

    # Slide 3: dataset structure
    slide = slides[2]
    clear_slide(slide)
    add_title(slide, "Dataset Structure", "The important signal is hidden in strings, groups, and small amounts of missing data")
    add_bullets_box(
        slide,
        0.7,
        1.25,
        4.9,
        2.7,
        "Core file facts",
        [
            f"train.csv: {stats['train_n']:,} rows x 14 columns",
            f"test.csv: {stats['test_n']:,} rows x 13 columns",
            f"75.99% of training rows are complete; 21.48% miss only one value.",
        ],
        ACCENT_BLUE,
    )
    add_bullets_box(
        slide,
        0.7,
        4.15,
        4.9,
        2.1,
        "Hidden structure",
        [
            "PassengerId encodes travel group and within-group position.",
            "Cabin encodes deck, cabin number, and ship side.",
            "Name can reveal surnames and family structure.",
        ],
        ACCENT_ORANGE,
    )
    slide.shapes.add_picture(str(charts["missingness"]), Inches(5.95), Inches(1.3), width=Inches(6.3))
    add_textbox(
        slide,
        6.15,
        5.95,
        5.95,
        0.42,
        "Blue = especially central for later parsing / imputation. Green = other partially missing features.",
        size=10,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Source: local train.csv/test.csv missingness audit")
    set_notes(
        slide,
        "This slide needs to do more than state file sizes; it needs to tell the audience what the variables actually represent in the story of the competition. The local files confirm that `train.csv` has 8,693 rows and 14 columns, while `test.csv` has 4,277 rows and 13 columns, with the only missing column in the test file being the response `Transported`. `Transported` is the binary outcome we are trying to predict: whether the passenger was sent to an alternate dimension. `PassengerId` is a unique passenger identifier, but it also embeds travel-group information because it is written as group underscore passenger number. `HomePlanet` is the passenger's planet of origin or residence. `CryoSleep` indicates whether the passenger spent the trip in suspended animation. `Cabin` is the passenger's cabin location and really means deck, cabin number, and ship side all at once. `Destination` is the exoplanet the passenger was traveling to. `Age` is literal age, and `VIP` indicates whether the passenger paid for premium service. The five spending columns, `RoomService`, `FoodCourt`, `ShoppingMall`, `Spa`, and `VRDeck`, are amounts billed to different amenities on the ship. `Name` is the passenger's full name and can sometimes reveal family or surname relationships.\n\n"
        "The chart on the right is not a chart of prediction importance; it is simply a chart of missingness in the training set. The horizontal axis gives the percentage of rows that are blank for each feature, and the vertical axis lists the features. The visual point is that almost all the bars are small, usually around two to two and a half percent, so this is not a catastrophically incomplete dataset. But the missing values occur in variables that matter. The blue bars do not mean categorical and the green bars do not mean quantitative. They also do not mean important versus unimportant in some absolute sense. The blue bars are just a visual emphasis for features that become especially central in later preprocessing and imputation, namely `CryoSleep`, `Cabin`, and `HomePlanet`. Those fields are highlighted because later notebooks repeatedly use them when engineering structure or filling in missing values. The green bars are the rest of the partially missing variables, such as `Age`, `Destination`, `VIP`, and the spending columns. Those are still meaningful variables; they are simply not the three fields I chose to visually emphasize on this chart.\n\n"
        "A final idea to make explicit is that this dataset contains hidden structure in places where a beginner might see only text. `PassengerId` is not just an ID; it tells us who is traveling together. `Cabin` is not just a string; it encodes where the passenger stayed on the ship. `Name` is not just a label; it may identify surnames or family groupings. That is why the slide subtitle says the signal is hidden in strings, groups, and small amounts of missing data. The conclusion I would draw out loud is that the dataset is mostly complete, but the columns need to be interpreted carefully before modeling. Once we understand what the variables mean, it becomes clear why public solutions spend so much time parsing strings and imputing strategically instead of rushing straight to a model."
    )

    # Slide 4: metric and validation
    slide = slides[3]
    clear_slide(slide)
    add_title(slide, "Accuracy, Baselines, and Local Validation", "The metric is simple; the validation design is where serious workflows separate themselves")
    add_bullets_box(
        slide,
        0.7,
        1.3,
        4.0,
        2.1,
        "Two formulas worth saying out loud",
        [
            "Accuracy = (TP + TN) / (TP + TN + FP + FN)",
            "CV error averages validation mistakes across folds; CV accuracy = 1 - CV error",
        ],
        ACCENT_BLUE,
    )
    add_bullets_box(
        slide,
        0.7,
        3.7,
        4.0,
        2.25,
        "Public validation numbers",
        [
            "Patrick RF OOB score: 0.80605",
            "Patrick HistGradientBoosting 10-fold CV: 0.80778",
            "Maria tuned LightGBM best CV: 0.80663",
        ],
        ACCENT_GREEN,
    )
    slide.shapes.add_picture(str(charts["validation"]), Inches(5.0), Inches(1.25), width=Inches(7.5))
    add_footer(slide, "Sources: PatrickSVM notebook; Maria Aguilera notebook; Flatiron starter notebook")
    set_notes(
        slide,
        "This slide should make a clean distinction between the official metric and the workflow participants use to choose models.\n"
        "- Kaggle scores a single submission on hidden labels, but participants need local validation before they submit.\n"
        "- PatrickSVM reports both an OOB Random Forest score of 0.80605 and a 10-fold CV mean of 0.80778 for HistGradientBoosting.\n"
        "- Maria Aguilera compares many models under a common CV design and then tunes LightGBM to about 0.80663.\n"
        "- The ladder chart is useful because it shows how much of the gain comes from moving from a naive baseline to a serious engineered workflow.\n"
        "- If you want a critical-thinking sentence, say that random CV may be mildly optimistic here because PassengerId encodes groups, so related passengers can leak across folds."
    )

    # Slide 5: feature signals
    slide = slides[4]
    clear_slide(slide)
    add_title(slide, "Feature Engineering Signals in the Data", "Public feature choices are defensible because the training set itself points to them")
    slide.shapes.add_picture(str(charts["signals"]), Inches(0.72), Inches(1.28), width=Inches(7.2))
    add_bullets_box(
        slide,
        8.15,
        1.3,
        4.4,
        2.15,
        "Group structure",
        [
            f"{fmt_rate(stats['uniform_group_outcome'])} of travel groups have a single observed outcome.",
            f"Same HomePlanet within group: {fmt_rate(stats['same_homeplanet'])}.",
            f"Same CabinSide within group: {fmt_rate(stats['same_cabinside'])}.",
        ],
        ACCENT_GREEN,
    )
    add_bullets_box(
        slide,
        8.15,
        3.7,
        4.4,
        2.15,
        "Cabin and planet logic",
        [
            "Observed rows on decks A/B/C are all Europa.",
            "Observed rows on deck G are all Earth.",
            "That is why deck-based HomePlanet imputation is sensible.",
        ],
        ACCENT_ORANGE,
    )
    add_footer(slide, "Source: local train.csv feature-rate audit; patterns echoed in PatrickSVM and Maria Aguilera notebooks")
    set_notes(
        slide,
        "This slide should sound like an argument for why the public notebooks engineered the features they did. The four mini-charts on the left are all showing conditional transport rates, meaning the proportion of passengers who were transported within each subgroup. The top-left chart compares `CryoSleep=False` to `CryoSleep=True`, and it is one of the strongest raw relationships in the whole dataset: passengers in cryosleep are transported at a much higher rate. The top-right chart compares positive spend to zero spend. That matters because zero spend is not just another numeric threshold. It lines up closely with the competition story that cryosleep passengers are confined and do not use amenities. So when public notebooks create `TotalSpend` and `NoSpend`, they are not inventing arbitrary features; they are formalizing a pattern already visible in the data.\n\n"
        "The bottom-left chart compares port and starboard side, and the bottom-right chart compares decks. Those two charts are especially important because they show that location on the ship matters. Side alone matters somewhat, but deck matters even more, and the deck effect is not uniform. Some decks have much higher transport rates than others, which immediately suggests that splitting `Cabin` into deck and side is more useful than treating it as an opaque string. This is also where imputation enters the conversation. Once we see that deck is predictive, and once we know from the raw data that certain decks line up strongly with particular home planets, deck stops being only a predictive feature and starts becoming a clue we can use to fill in missing `HomePlanet` values sensibly.\n\n"
        "The boxes on the right are summarizing the same idea at the group level. The 'Group structure' box is there because `PassengerId` embeds travel groups, and those groups are not random bookkeeping. In the local training data, 87.18 percent of groups have a single observed transport outcome, which means people traveling together are often transported together. On top of that, when `HomePlanet` and `CabinSide` are observed within a group, they are essentially perfectly consistent. That gives a very strong justification for group-based imputation. The 'Cabin and planet logic' box highlights another important inference: decks A, B, and C are overwhelmingly Europa-associated in the observed data, while deck G is strongly associated with Earth. So if a passenger is on one of those decks and `HomePlanet` is missing, imputing from deck is not just guesswork; it is a data-supported rule. The sentence I would emphasize out loud is that good feature engineering here means decoding the structure the data already contain, especially group structure, ship location, and the logical link between spending and cryosleep."
    )

    # Slide 6: workflow
    slide = slides[5]
    clear_slide(slide)
    add_title(slide, "What Strong Public Workflows Actually Did", "The recurring process is disciplined and highly reproducible")
    steps = [
        ("1. Audit", "Inspect distributions, blanks, and target balance"),
        ("2. Decode", "Parse PassengerId, Cabin, and sometimes surnames"),
        ("3. Engineer", "Create GroupSize, TotalSpend, NoSpend, deck/side, cabin regions"),
        ("4. Impute", "Use cryosleep, deck, and group logic rather than blind means"),
        ("5. Validate", "Compare on OOB or K-fold CV before submission"),
        ("6. Refine", "Tune boosting models and sometimes ensemble the strongest ones"),
    ]
    x_positions = [0.7, 2.85, 5.0, 7.15, 9.3, 11.45]
    for (label, text), x in zip(steps, x_positions):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(1.55), Inches(2.55)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL_DARK
        shape.line.color.rgb = ACCENT_BLUE if x < 7 else ACCENT_GREEN
        shape.line.width = Pt(1.0)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            set_run_style(run, 14, ACCENT_BLUE if x < 7 else ACCENT_GREEN, bold=True)
        p = tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            set_run_style(run, 11, WHITE)
    add_textbox(
        slide,
        0.8,
        5.15,
        12.0,
        0.7,
        "Across PatrickSVM, Maria Aguilera, Amir Fares, and Samyak Raj Bayar, the repeated theme is the same: "
        "feature engineering and logic-based preprocessing happen before serious model comparison.",
        size=16,
    )
    add_footer(slide, "Sources: PatrickSVM; Maria Aguilera; Amir Fares; Samyak Raj Bayar")
    set_notes(
        slide,
        "This is the workflow slide, so use it to move from raw data facts to public practice.\n"
        "- PatrickSVM starts with extensive EDA and case-based imputation before fitting Random Forest and HistGradientBoosting.\n"
        "- Maria Aguilera follows the same general order: decode PassengerId and Cabin, create no-spend and family-related features, impute with deck and group logic, then compare models and tune them.\n"
        "- Amir Fares is less formal in exposition but still uses the same structure: preprocess first, try many models second, ensemble last.\n"
        "- Samyak Raj Bayar's baseline is useful because it shows that even a short public script still parses Cabin and GroupSize before using HistGradientBoosting.\n"
        "- The phrase to emphasize is that these are workflow improvements, not just model swaps."
    )

    # Slide 7: model comparison
    slide = slides[6]
    clear_slide(slide)
    add_title(slide, "Model Comparison from Public Notebooks", "Good preprocessing puts many models near 0.80; boosting refines the last mile")
    slide.shapes.add_picture(str(charts["models"]), Inches(0.8), Inches(1.25), width=Inches(7.6))
    add_bullets_box(
        slide,
        8.65,
        1.35,
        3.9,
        1.85,
        "What the chart means",
        [
            "Engineered logistic is already around 0.795.",
            "Untuned tree ensembles cluster near 0.80.",
            "Tuning moves the best boosted models a bit higher.",
        ],
        ACCENT_BLUE,
    )
    add_bullets_box(
        slide,
        8.65,
        3.55,
        3.9,
        2.05,
        "The nuanced conclusion",
        [
            "Most of the gain comes from preprocessing, not from a magical model.",
            "Boosted trees still win the strongest public comparisons.",
            "Deep learning appears in experiments, but not as the dominant public story.",
        ],
        ACCENT_GREEN,
    )
    add_footer(slide, "Sources: Maria Aguilera notebook; PatrickSVM notebook; Amir Fares repo")
    set_notes(
        slide,
        "This slide is where the presentation should avoid oversimplifying the model story.\n"
        "- Maria Aguilera's results show that once preprocessing is strong, several models live in a tight band around 0.80.\n"
        "- That means feature engineering explains a large share of the final performance.\n"
        "- The nonlinear tree models still matter because they capture deck-by-side and cryosleep-by-spend interactions automatically.\n"
        "- A useful quantitative sentence is that tuned LightGBM at 0.8066 is only about 1.15 percentage points above the engineered logistic baseline at 0.7951.\n"
        "- So the public evidence favors boosted trees, but it also says the preprocessing pipeline is doing most of the heavy lifting."
    )

    # Slide 8: LightGBM deep dive
    slide = slides[7]
    clear_slide(slide)
    add_title(slide, "LightGBM in Detail", "Maria Aguilera's tuned LightGBM shows why histogram boosting is so strong on tabular data")
    add_diagram_box(
        slide,
        0.8,
        1.45,
        2.1,
        1.45,
        "1. Bin Features",
        "Continuous variables are quantized into histogram bins, so LightGBM scans bins instead of every raw split value.",
        ACCENT_BLUE,
        title_size=13,
        body_size=10,
    )
    add_diagram_box(
        slide,
        3.08,
        1.45,
        2.1,
        1.45,
        "2. Grow Best Leaf",
        "LightGBM grows trees leaf-wise: it expands the leaf with the biggest loss reduction, not every level evenly.",
        ACCENT_GREEN,
        title_size=13,
        body_size=10,
    )
    add_diagram_box(
        slide,
        5.36,
        1.45,
        2.1,
        1.45,
        "3. Update Ensemble",
        "Add the new tree to the ensemble and repeat for many rounds; Maria's tuned model used 400 estimators.",
        ACCENT_ORANGE,
        title_size=13,
        body_size=10,
    )
    add_arrow(slide, 2.87, 1.95, 0.16, 0.28, direction="right", color=ACCENT_BLUE)
    add_arrow(slide, 5.15, 1.95, 0.16, 0.28, direction="right", color=ACCENT_GREEN)
    add_arrow(slide, 4.02, 3.12, 0.24, 0.34, direction="down", color=WHITE)
    add_diagram_box(
        slide,
        1.15,
        3.45,
        5.95,
        1.6,
        "Why LightGBM fits this competition",
        "Histogram speed plus leaf-wise growth make LightGBM strong on mixed tabular data. It can learn interactions like CryoSleep by spend or Deck by Side while remaining efficient enough to tune aggressively.",
        ACCENT_TEAL,
        title_size=14,
        body_size=10,
    )
    add_textbox(
        slide,
        2.05,
        5.18,
        4.2,
        0.35,
        "Maria Aguilera tuned LightGBM CV = 0.80663",
        size=11,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_bullets_box(
        slide,
        8.2,
        1.3,
        4.3,
        2.35,
        "Boosting / Sampling Controls",
        [
            "n_estimators=400: many boosted trees, giving the ensemble time to refine errors gradually.",
            "subsample=0.9: each boosting round sees 90% of rows, adding regularization through bagging.",
            "subsample_freq=20: row subsampling is applied periodically rather than on every single tree.",
            "colsample_bytree=0.8: each tree sees 80% of features, which reduces correlation and overfitting.",
        ],
        ACCENT_BLUE,
        title_size=17,
        bullet_size=12,
    )
    add_bullets_box(
        slide,
        8.2,
        3.95,
        4.3,
        2.3,
        "Tree / Data Controls",
        [
            "num_leaves=100: the main LightGBM complexity control under leaf-wise growth.",
            "max_depth=20: caps how deep leaf-wise growth is allowed to go.",
            "min_split_gain=0.4: a split must improve the objective enough to be worth taking.",
            "reg_alpha=1.3 and reg_lambda=1.1: L1 and L2 penalties regularize leaf weights.",
        ],
        ACCENT_GREEN,
        title_size=17,
        bullet_size=12,
    )
    add_textbox(
        slide,
        8.28,
        6.35,
        4.15,
        0.46,
        "This slide sets up HistGB: same histogram idea next, but inside scikit-learn rather than an external library.",
        size=11,
        color=MUTED,
    )
    add_footer(slide, "Sources: Maria Aguilera notebook; official LightGBM documentation")
    set_notes(
        slide,
        "This slide should introduce LightGBM as the tuned, external-library version of histogram-based gradient boosting before we move to HistGradientBoosting on the next slide. The left-side diagram is the most important visual. The first box explains the core efficiency trick: LightGBM bins continuous values into histograms and scans those bins for splits, which is much cheaper than exact split search on every unique value. The second box explains the second major idea: unlike level-wise boosting, LightGBM grows trees leaf-wise, meaning it keeps splitting whichever current leaf gives the largest reduction in loss. That is one reason LightGBM can be so strong on tabular problems, but it is also why parameters like `num_leaves` and `max_depth` matter so much. The bottom box is the payoff: Spaceship Titanic has mixed tabular features and strong interactions, so LightGBM can model those efficiently.\n\n"
        "The right side should be read as Maria Aguilera's tuned implementation choices. `n_estimators=400` means the ensemble is quite large, with many opportunities to reduce residual error. `subsample=0.9` and `subsample_freq=20` are row-sampling controls that regularize the fit, while `colsample_bytree=0.8` does the same on the feature side by limiting each tree to 80 percent of the predictors. Then the lower box is about structural complexity and regularization. `num_leaves=100` is arguably the main LightGBM complexity control because LightGBM grows leaf-wise. `max_depth=20` puts a hard ceiling on that growth. `min_split_gain=0.4` says a split must improve the objective enough to be worthwhile, and `reg_alpha` plus `reg_lambda` regularize the learned leaf weights.\n\n"
        "The transition sentence to use out loud is this: LightGBM shows why histogram boosting is attractive in the first place, and HistGradientBoosting is the natural next slide because it keeps the histogram idea but moves back into native scikit-learn tooling. Maria's tuned LightGBM reaches about 0.80663 in cross-validation, so it gives us a strong benchmark before we compare it to Patrick's HistGB setup."
    )

    # Slide 9: HistGradientBoosting deep dive
    slide = slides[8]
    clear_slide(slide)
    add_title(slide, "HistGradientBoostingClassifier in Detail", "HistGB keeps the histogram idea, but does it inside native scikit-learn")
    add_diagram_box(
        slide,
        0.8,
        1.45,
        2.1,
        1.45,
        "1. Bin Features",
        "Like LightGBM, HistGB compresses numeric values into bins, which cuts split-search cost and memory use.",
        ACCENT_BLUE,
        title_size=13,
        body_size=10,
    )
    add_diagram_box(
        slide,
        3.08,
        1.45,
        2.1,
        1.45,
        "2. Fit Weak Tree",
        "At each round, fit a tree to the negative gradient of binary log-loss.",
        ACCENT_GREEN,
        title_size=13,
        body_size=10,
    )
    add_diagram_box(
        slide,
        5.36,
        1.45,
        2.1,
        1.45,
        "3. Update Ensemble",
        "Add that tree back with a small step size, then repeat for 120 boosting rounds.",
        ACCENT_ORANGE,
        title_size=13,
        body_size=10,
    )
    add_arrow(slide, 2.87, 1.95, 0.16, 0.28, direction="right", color=ACCENT_BLUE)
    add_arrow(slide, 5.15, 1.95, 0.16, 0.28, direction="right", color=ACCENT_GREEN)
    add_arrow(slide, 4.02, 3.12, 0.24, 0.34, direction="down", color=WHITE)
    add_diagram_box(
        slide,
        1.15,
        3.45,
        5.95,
        1.6,
        "Why Hist matters here",
        "The strength is not only boosting. Histogram binning makes repeated split search much cheaper, so 120 rounds of interaction-rich tree fitting stay practical even with mixed parsed features.",
        ACCENT_TEAL,
        title_size=14,
        body_size=10,
    )
    add_textbox(
        slide,
        2.05,
        5.18,
        4.2,
        0.35,
        "PatrickSVM 10-fold CV mean = 0.80778",
        size=11,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_bullets_box(
        slide,
        8.2,
        1.3,
        4.3,
        2.35,
        "Boosting Controls",
        [
            "loss='log_loss': optimized for binary classification, not regression.",
            "max_iter=120: 120 boosting stages, so the ensemble keeps correcting residual mistakes.",
            "learning_rate=0.075: each stage only nudges the fit, which stabilizes boosting.",
        ],
        ACCENT_BLUE,
        title_size=17,
        bullet_size=12,
    )
    add_bullets_box(
        slide,
        8.2,
        3.95,
        4.3,
        2.3,
        "Tree / Data Controls",
        [
            "max_depth=10 and max_leaf_nodes=16: trees can learn interactions, but each stage is still bounded.",
            "min_samples_leaf=6: tiny terminal nodes are disallowed, which regularizes the fit.",
            "categorical_features=mask: lets HistGB treat parsed fields like Deck, Side, GroupSize, and CryoSleep as categorical.",
            "random_state=123 fixes reproducibility; l2_regularization=0 means shrinkage and tree shape do most of the regularization.",
        ],
        ACCENT_GREEN,
        title_size=17,
        bullet_size=12,
    )
    add_textbox(
        slide,
        8.28,
        6.35,
        4.15,
        0.46,
        "Patrick also used one_hot=False and log-transformed expenses before fitting HistGB.",
        size=11,
        color=MUTED,
    )
    add_footer(slide, "Sources: PatrickSVM notebook; official scikit-learn HistGradientBoosting docs")
    set_notes(
        slide,
        "This slide should now feel like a natural continuation of the LightGBM slide. The transition is that both models use histogram-based split finding, but HistGradientBoosting is the native scikit-learn version rather than an external boosting library. That matters for a class presentation because it shows that the histogram trick is bigger than one package. The first box should be explained very specifically: the strength of 'Hist' is that numeric values are compressed into bins, and split search is performed on those bins. That reduces computational cost and memory usage relative to exact split search, which is why repeated boosting rounds remain practical.\n\n"
        "The rest of the left diagram is ordinary boosting logic, but the histogram point is what differentiates this from generic gradient boosting. After the features are binned, HistGB fits a weak tree to the negative gradient of log-loss, adds that tree back to the ensemble with a small step, and repeats the process. The lower box should be read as a claim about why histogram boosting is useful here specifically: once we parse `Cabin`, `PassengerId`, and spending features, we have a mixed tabular dataset with interactions and skew, and histogram boosting is a computationally efficient way to fit many rounds of trees without hand-writing interaction terms.\n\n"
        "On the right, the parameters should be interpreted as PatrickSVM's concrete implementation choices. `max_iter=120` and `learning_rate=0.075` are the classic boosting capacity-versus-shrinkage pair. `max_depth=10`, `max_leaf_nodes=16`, and `min_samples_leaf=6` keep the individual trees expressive enough to learn interactions but regularized enough not to explode in complexity. `categorical_features=mask` is especially important, because after preprocessing Patrick can tell HistGB which fields should be treated as categorical rather than purely numeric. The small footer sentence about `one_hot=False` and log-transformed spending matters for the same reason it mattered on the LightGBM slide: preprocessing and model design are working together. The conclusion to land on is that HistGB is not just 'another boosting model'; it is the scikit-learn-native version of the same histogram logic that makes LightGBM strong."
    )

    # Slide 10: ISLR and trade-offs
    slide = slides[9]
    clear_slide(slide)
    add_title(slide, "ISLR Connections and Trade-offs", "Why this competition sits at the boundary of classical statistical learning and practical ML")
    add_bullets_box(
        slide,
        0.72,
        1.3,
        3.0,
        2.2,
        "Cross-validation",
        [
            "Use local K-fold CV because test labels are hidden.",
            "Patrick HGB 10-fold CV mean: 0.80778.",
            "Model selection here is a real CV problem, not a leaderboard guess.",
        ],
        ACCENT_BLUE,
    )
    add_bullets_box(
        slide,
        3.98,
        1.3,
        3.0,
        2.2,
        "Logistic regression",
        [
            "Useful baseline once features are engineered and encoded.",
            "Maria engineered logistic CV mean: 0.7951.",
            "Additive log-odds structure misses interactions unless we add them by hand.",
        ],
        ACCENT_GREEN,
    )
    add_bullets_box(
        slide,
        7.24,
        1.3,
        2.95,
        2.2,
        "Ridge / Lasso",
        [
            "Natural baseline when one-hot features and interactions make p grow.",
            "Shrink unstable coefficients and control collinearity.",
            "Still need explicit interaction design.",
        ],
        ACCENT_ORANGE,
    )
    add_bullets_box(
        slide,
        10.45,
        1.3,
        2.1,
        2.2,
        "Bias-variance",
        [
            "RF is stable.",
            "Boosting is slightly better but more tune-sensitive.",
            "Gains are real, but not huge after good preprocessing.",
        ],
        ACCENT_TEAL,
    )
    add_textbox(
        slide,
        0.9,
        4.15,
        11.6,
        1.35,
        "Useful formulas to reference verbally:\n"
        "P(Y=1|X=x) = exp(beta_0 + x^T beta) / (1 + exp(beta_0 + x^T beta))\n"
        "Ridge logistic: -log L(beta) + lambda * sum beta_j^2\n"
        "Lasso logistic: -log L(beta) + lambda * sum |beta_j|",
        size=15,
    )
    add_footer(slide, "Sources: Maria Aguilera notebook; PatrickSVM notebook; course concepts from ISLR")
    set_notes(
        slide,
        "This is the slide where I would deliberately shift into course language and show that the competition can be understood through ISLR rather than only through Kaggle jargon. The first box is cross-validation, and the point there is that the hidden test set forces participants to estimate test error indirectly. That is exactly the model-selection problem we study in class. Patrick's HistGradientBoosting result of about 0.80778 from 10-fold cross-validation is a concrete example of using resampling to compare candidates before submission, instead of trusting the public leaderboard alone. If I wanted one sentence that sounds very course-aligned, it would be: cross-validation is how participants estimate expected generalization performance when the real test labels are unavailable.\n\n"
        "The logistic-regression box is useful because it keeps the competition connected to classical classification. I would say that logistic regression is not weak here because it is old; it is only limited because its default form is additive on the log-odds scale. Once Maria Aguilera engineers the features well, logistic regression reaches roughly 0.795 in cross-validation, which is already strong. The issue is that Spaceship Titanic contains interactions that a plain linear log-odds model does not automatically recover. CryoSleep and no-spend reinforce each other. Deck and side matter together, not independently. Group structure changes how some fields should be interpreted. So this is a good example of a dataset where classical models remain competitive, but only after the analyst puts substantial work into encoding the right structure.\n\n"
        "The ridge and lasso box is there to connect the competition to the regularization material in class. If we one-hot encode deck, side, destination, home planet, VIP, cryosleep, and then start adding interactions, the design matrix becomes wider and more collinear. Ridge would stabilize coefficient estimates by shrinking them continuously, lasso could perform some automatic variable selection by shrinking some coefficients to zero, and elastic net would split the difference. So even though most public notebooks eventually prefer boosted trees, penalized logistic regression would still be a principled and very course-appropriate baseline. That is an important nuance: the best Kaggle approach is not always the only statistically sensible approach.\n\n"
        "The bias-variance trade-off is the final interpretive piece. Random Forest tends to be a stable, lower-risk ensemble because it averages many decorrelated trees; HistGB and LightGBM are more flexible and can usually do a bit better, but they also need more tuning discipline. The key phrase here is 'a bit better.' Once preprocessing is good, the gap between engineered logistic regression and the best boosted methods is measured in about one percentage point, not in a huge leap. That is a very good ISLR lesson: flexible models can reduce bias, but after a strong feature-engineering pipeline the marginal gains become smaller, and the cost is increased complexity. So the trade-off is real rather than rhetorical."
    )

    # Slide 11: best practices
    slide = slides[10]
    clear_slide(slide)
    add_title(slide, "Best Practices and Reflection", "Transferable lessons that matter outside this specific competition")
    practices = [
        "Decode structured strings before dropping them.",
        "Impute with domain logic, not just global means or modes.",
        "Separate feature gains from model gains.",
        "Compare models on the same validation design.",
        "Report limitations, leakage risks, and leaderboard instability honestly.",
    ]
    for idx, practice in enumerate(practices, start=1):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.2 + (idx - 1) * 1.02), Inches(11.6), Inches(0.72)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL_DARK if idx % 2 else PANEL
        shape.line.color.rgb = ACCENT_GREEN if idx % 2 else ACCENT_BLUE
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(12)
        p = tf.paragraphs[0]
        p.text = f"{idx}. {practice}"
        for run in p.runs:
            set_run_style(run, 18, WHITE, bold=False)
    add_textbox(
        slide,
        1.0,
        6.48,
        11.3,
        0.35,
        "Bottom line: Spaceship Titanic rewards statistical thinking as much as model choice.",
        size=18,
        color=ACCENT_GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Sources: local CSV audit plus public workflow synthesis across multiple notebooks")
    set_notes(
        slide,
        "I would use this slide as the synthesis of the whole presentation, not as a list of generic advice. The first best practice, decoding structured strings before dropping them, comes directly from the competition itself. `PassengerId` contains travel-group information, `Cabin` contains deck and side, and even `Name` can reveal family structure. A weaker workflow would throw those away because they look like text. A stronger workflow asks whether the text is actually compressed structure. That is one of the most transferable lessons from this project: before discarding a variable that looks messy, ask what real-world process generated it.\n\n"
        "The second best practice is logic-based imputation. In this dataset, missing values are not random holes in a spreadsheet. They sit inside a network of relationships. CryoSleep is tied to spending. Group members usually share home planet and cabin side. Certain decks align strongly with certain planets. So the best public notebooks do not simply plug in global means and modes. They use the relationships among variables to recover plausible missing values. That is good statistical practice because it uses the data-generating structure rather than ignoring it.\n\n"
        "The third and fourth practices are about disciplined modeling. Separating feature gains from model gains matters because it keeps us honest about where performance improvements are coming from. A one-feature baseline around 0.63, an engineered logistic baseline around 0.795, and boosted-tree models around 0.806 to 0.808 tell a very different story than simply saying 'boosting wins.' Likewise, comparing models under the same validation design matters because otherwise we are mixing together preprocessing effects, split effects, and model effects. PatrickSVM, Maria Aguilera, and other strong public workflows are credible not just because their scores are good, but because they validate locally and document what changed.\n\n"
        "The final reflection I would give is that Spaceship Titanic is a very good teaching competition because it rewards statistical thinking at every step. The strongest workflows are not magical. They are careful. They inspect the data, decode hidden structure, impute logically, validate consistently, and only then tune models. That is why this project is more valuable than just memorizing which algorithm got the highest score. It shows that good machine learning practice is mostly good reasoning, applied carefully and reproducibly."
    )

    # Slide 12: sources
    slide = slides[11]
    clear_slide(slide)
    add_title(slide, "Selected Sources", "Primary references used to build the deck and speaker notes")
    add_textbox(
        slide,
        0.9,
        1.35,
        11.5,
        4.9,
        "1. Official Kaggle competition page: Spaceship Titanic\n"
        "2. Class competition brief and rubric\n"
        "3. PatrickSVM README and public notebook\n"
        "4. Maria Aguilera exported notebook\n"
        "5. Amir Fares repo and weighted-ensemble notebook link\n"
        "6. Flatiron starter notebook for a simple logistic baseline\n"
        "7. Samyak Raj Bayar public HistGradientBoosting baseline\n"
        "8. Official scikit-learn HistGradientBoosting documentation\n"
        "9. Official LightGBM documentation\n\n"
        "All dataset counts and conditional rates shown in the slides come from the attached train.csv/test.csv.",
        size=18,
    )
    add_footer(slide, "Deck generated programmatically from kaggle.md and the attached CSV files")
    set_notes(
        slide,
        "Close by noting that the numbers in the slides are not generic internet facts.\n"
        "- The competition and rubric framing came from the assignment context.\n"
        "- The descriptive statistics came from the attached local CSVs.\n"
        "- The workflow and model-comparison claims came from public notebooks and repositories that are specific to Spaceship Titanic.\n"
        "- The LightGBM slide also uses the official LightGBM documentation for the histogram and leaf-wise interpretations.\n"
        "- The HistGradientBoosting explainer slide also uses the official scikit-learn documentation for the histogram-binning and categorical-feature details.\n"
        "- If someone asks where a number came from, the notes and this slide give a clean answer."
    )

    prs.save(PPTX_PATH)


if __name__ == "__main__":
    build_presentation()
